#!/usr/bin/env python3
"""
make_video.py — slide deck → narrated lecture video pipeline.

Pipeline (matches the recommended free pipeline in
``docs/research/teaching-video-tools-2026.md``):

    .pptx  --libreoffice headless-->  .pdf
           --pdf2image (pdftoppm)-->  one PNG per slide (1920x1080)
    .pptx  --python-pptx        -->  speaker-notes script per slide
    notes  --gTTS / pyttsx3     -->  per-slide narration (mp3/wav)
    PNG+WAV--ffmpeg             -->  per-slide subclip
    subclips--ffmpeg xfade      -->  final .mp4
    ffprobe                     -->  duration check

The script is intentionally tolerant: any missing optional binary (libreoffice,
ffmpeg, ffprobe) or library (gTTS, pyttsx3) prints a clear actionable error
instead of crashing with a traceback. Slide PNGs are cached under a
deterministic temp directory so re-runs are fast.

Usage:
    python3 tools/make_video.py --input book/decks/week-01.pptx \\
                                --output videos/camp-week-01.mp4

Optional flags:
    --tts {auto,gtts,pyttsx3}   force a TTS backend (default: auto)
    --dpi N                     pdftoppm DPI (default: 150)
    --crossfade SECS            transition duration (default: 0.3)
    --pre-pad SECS              silence before first slide (default: 1.0)
    --post-pad SECS             silence after last slide (default: 1.5)
    --cache-dir DIR             reuse this dir for intermediate artifacts
    --keep-temp                 don't delete the cache dir on exit
    --lang CODE                 gTTS language (default: en)
    --no-crossfade              concat without xfade (faster, less polished)
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import shutil
import subprocess
import sys
import tempfile
import textwrap
import time
from dataclasses import dataclass
from pathlib import Path


# ---------------------------------------------------------------------------
# Small helpers
# ---------------------------------------------------------------------------

def _eprint(msg: str) -> None:
    print(msg, file=sys.stderr, flush=True)


def _log(msg: str) -> None:
    print(f"[make_video] {msg}", flush=True)


def _die(msg: str, code: int = 1) -> "NoReturn":  # type: ignore[name-defined]
    _eprint(f"[make_video] ERROR: {msg}")
    sys.exit(code)


def _which(*candidates: str) -> str | None:
    for c in candidates:
        p = shutil.which(c)
        if p:
            return p
    return None


def _run(cmd: list[str], *, check: bool = True, capture: bool = False) -> subprocess.CompletedProcess:
    """Run a subprocess with helpful error messages."""
    try:
        return subprocess.run(
            cmd,
            check=check,
            stdout=subprocess.PIPE if capture else None,
            stderr=subprocess.PIPE if capture else None,
            text=True,
        )
    except FileNotFoundError as e:
        _die(f"required binary not found: {cmd[0]} ({e})")
    except subprocess.CalledProcessError as e:
        tail = ""
        if capture and e.stderr:
            tail = "\n--- stderr ---\n" + e.stderr.strip()[-2000:]
        _die(f"command failed (exit {e.returncode}): {' '.join(cmd)}{tail}")


# ---------------------------------------------------------------------------
# Step 1 — slide images
# ---------------------------------------------------------------------------

def pptx_to_pdf(pptx: Path, out_dir: Path) -> Path:
    """Convert .pptx to .pdf via libreoffice headless."""
    soffice = _which("libreoffice", "soffice")
    if not soffice:
        _die(
            "libreoffice / soffice not found on PATH. Install with:\n"
            "    sudo apt-get install -y libreoffice-nogui\n"
            "or on macOS:  brew install --cask libreoffice"
        )
    out_dir.mkdir(parents=True, exist_ok=True)
    _log(f"libreoffice: {pptx.name} -> PDF")
    _run(
        [
            soffice,
            "--headless",
            "--invisible",
            "--nologo",
            "--nodefault",
            "--norestore",
            "--convert-to",
            "pdf",
            "--outdir",
            str(out_dir),
            str(pptx),
        ],
        capture=True,
    )
    pdf = out_dir / (pptx.stem + ".pdf")
    if not pdf.exists():
        _die(f"libreoffice did not produce a PDF at {pdf}")
    return pdf


def pdf_to_pngs(pdf: Path, out_dir: Path, dpi: int = 150, width: int = 1920, height: int = 1080) -> list[Path]:
    """Convert PDF to one PNG per slide using pdf2image, then resize to 1920x1080."""
    try:
        from pdf2image import convert_from_path  # type: ignore
    except ImportError:
        _die(
            "pdf2image not installed. Install with:\n"
            "    pip install pdf2image\n"
            "and the poppler binaries:\n"
            "    sudo apt-get install -y poppler-utils"
        )

    if not _which("pdftoppm"):
        _die(
            "poppler-utils (pdftoppm) not found on PATH. Install with:\n"
            "    sudo apt-get install -y poppler-utils"
        )

    out_dir.mkdir(parents=True, exist_ok=True)
    _log(f"pdf2image: rendering {pdf.name} at {dpi} DPI")
    images = convert_from_path(str(pdf), dpi=dpi)

    try:
        from PIL import Image  # type: ignore
    except ImportError:
        _die("Pillow not installed. Install with:  pip install Pillow")

    paths: list[Path] = []
    for i, im in enumerate(images, start=1):
        # Resize/letterbox to 1920x1080 (preserve aspect ratio, pad with black).
        canvas = Image.new("RGB", (width, height), (0, 0, 0))
        ratio = min(width / im.width, height / im.height)
        new_w, new_h = int(im.width * ratio), int(im.height * ratio)
        resized = im.resize((new_w, new_h), Image.LANCZOS)
        canvas.paste(resized, ((width - new_w) // 2, (height - new_h) // 2))
        png = out_dir / f"slide-{i:03d}.png"
        canvas.save(png, "PNG")
        paths.append(png)
    _log(f"pdf2image: wrote {len(paths)} PNG(s) to {out_dir}")
    return paths


# ---------------------------------------------------------------------------
# Step 2 — speaker notes
# ---------------------------------------------------------------------------

def extract_notes(pptx: Path) -> list[str]:
    """Read speaker notes from every slide. Empty notes become a 1-word placeholder
    so the TTS step never produces a zero-length clip."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        _die("python-pptx not installed. Install with:  pip install python-pptx")
    prs = Presentation(str(pptx))
    notes: list[str] = []
    for slide in prs.slides:
        text = ""
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            if tf is not None:
                text = (tf.text or "").strip()
        if not text:
            text = "(no narration provided for this slide.)"
        notes.append(text)
    _log(f"python-pptx: extracted notes for {len(notes)} slide(s)")
    return notes


# ---------------------------------------------------------------------------
# Step 3 — TTS
# ---------------------------------------------------------------------------

@dataclass
class TTSResult:
    backend: str
    files: list[Path]


def _tts_gtts(scripts: list[str], out_dir: Path, lang: str) -> list[Path]:
    from gtts import gTTS  # type: ignore
    paths: list[Path] = []
    for i, text in enumerate(scripts, start=1):
        mp3 = out_dir / f"narration-{i:03d}.mp3"
        if mp3.exists() and mp3.stat().st_size > 0:
            paths.append(mp3)
            continue
        _log(f"gTTS: synth slide {i}/{len(scripts)} ({len(text)} chars)")
        tts = gTTS(text=text, lang=lang)
        tts.save(str(mp3))
        paths.append(mp3)
    return paths


def _tts_pyttsx3(scripts: list[str], out_dir: Path) -> list[Path]:
    import pyttsx3  # type: ignore
    engine = pyttsx3.init()
    paths: list[Path] = []
    for i, text in enumerate(scripts, start=1):
        wav = out_dir / f"narration-{i:03d}.wav"
        if wav.exists() and wav.stat().st_size > 0:
            paths.append(wav)
            continue
        _log(f"pyttsx3: synth slide {i}/{len(scripts)} ({len(text)} chars)")
        engine.save_to_file(text, str(wav))
        engine.runAndWait()
        paths.append(wav)
    return paths


def synth_tts(scripts: list[str], out_dir: Path, *, choice: str = "auto", lang: str = "en") -> TTSResult:
    out_dir.mkdir(parents=True, exist_ok=True)
    tried: list[str] = []

    def _try_gtts() -> list[Path] | None:
        try:
            import gtts  # noqa: F401
        except ImportError:
            tried.append("gTTS (not installed)")
            return None
        try:
            return _tts_gtts(scripts, out_dir, lang)
        except Exception as e:  # network errors, rate-limit, etc.
            tried.append(f"gTTS (failed: {e.__class__.__name__}: {e})")
            return None

    def _try_pyttsx3() -> list[Path] | None:
        try:
            import pyttsx3  # noqa: F401
        except ImportError:
            tried.append("pyttsx3 (not installed)")
            return None
        try:
            return _tts_pyttsx3(scripts, out_dir)
        except Exception as e:
            tried.append(f"pyttsx3 (failed: {e.__class__.__name__}: {e})")
            return None

    if choice in ("auto", "gtts"):
        files = _try_gtts()
        if files is not None:
            return TTSResult("gTTS", files)
        if choice == "gtts":
            _die(f"gTTS backend requested but unavailable. Tried: {tried}")

    if choice in ("auto", "pyttsx3"):
        files = _try_pyttsx3()
        if files is not None:
            return TTSResult("pyttsx3", files)
        if choice == "pyttsx3":
            _die(f"pyttsx3 backend requested but unavailable. Tried: {tried}")

    _die(
        "no TTS backend available. Install one of:\n"
        "    pip install gTTS         # network, decent quality\n"
        "    pip install pyttsx3       # offline, robotic\n"
        f"Backends tried: {tried}"
    )


# ---------------------------------------------------------------------------
# Step 4 — ffmpeg assembly
# ---------------------------------------------------------------------------

def _need_ffmpeg() -> str:
    ff = _which("ffmpeg")
    if not ff:
        _die(
            "ffmpeg not found on PATH. Install with:\n"
            "    sudo apt-get install -y ffmpeg\n"
            "or:  brew install ffmpeg"
        )
    return ff


def _need_ffprobe() -> str:
    fp = _which("ffprobe")
    if not fp:
        _die(
            "ffprobe not found on PATH. It ships with ffmpeg — install with:\n"
            "    sudo apt-get install -y ffmpeg"
        )
    return fp


def audio_duration(path: Path) -> float:
    ffprobe = _need_ffprobe()
    out = _run(
        [
            ffprobe,
            "-v",
            "error",
            "-show_entries",
            "format=duration",
            "-of",
            "default=noprint_wrappers=1:nokey=1",
            str(path),
        ],
        capture=True,
    )
    return float(out.stdout.strip())


def make_slide_clip(
    png: Path,
    audio: Path,
    out_mp4: Path,
    duration: float,
    *,
    pre_pad: float = 0.0,
    post_pad: float = 0.0,
) -> Path:
    """Render one slide as a still-image video matched to its narration."""
    ffmpeg = _need_ffmpeg()
    total = duration + pre_pad + post_pad

    # Build the audio filter: pad with silence on either side if requested.
    a_filter_parts: list[str] = []
    if pre_pad > 0:
        a_filter_parts.append(f"adelay={int(pre_pad*1000)}|{int(pre_pad*1000)}")
    if post_pad > 0:
        a_filter_parts.append(f"apad=pad_dur={post_pad}")
    a_filter = ",".join(a_filter_parts) if a_filter_parts else "anull"

    cmd = [
        ffmpeg,
        "-y",
        "-loglevel",
        "error",
        "-loop",
        "1",
        "-i",
        str(png),
        "-i",
        str(audio),
        "-filter_complex",
        f"[0:v]scale=1920:1080:force_original_aspect_ratio=decrease,"
        f"pad=1920:1080:(ow-iw)/2:(oh-ih)/2,setsar=1,fps=30[v];"
        f"[1:a]{a_filter}[a]",
        "-map",
        "[v]",
        "-map",
        "[a]",
        "-c:v",
        "libx264",
        "-pix_fmt",
        "yuv420p",
        "-preset",
        "medium",
        "-crf",
        "20",
        "-c:a",
        "aac",
        "-b:a",
        "192k",
        "-shortest",
        "-t",
        f"{total:.3f}",
        str(out_mp4),
    ]
    _run(cmd, capture=True)
    return out_mp4


def concat_clips(clips: list[Path], out_mp4: Path, *, crossfade: float = 0.3) -> None:
    """Stitch per-slide clips. With crossfade>0 use xfade/acrossfade; else concat."""
    ffmpeg = _need_ffmpeg()
    out_mp4.parent.mkdir(parents=True, exist_ok=True)

    if len(clips) == 1:
        shutil.copyfile(clips[0], out_mp4)
        return

    if crossfade <= 0:
        # Demuxer-based concat (re-encode for safety since clips share params).
        list_file = out_mp4.parent / (out_mp4.stem + ".concat.txt")
        list_file.write_text("".join(f"file '{c.resolve()}'\n" for c in clips))
        _run(
            [
                ffmpeg,
                "-y",
                "-loglevel",
                "error",
                "-f",
                "concat",
                "-safe",
                "0",
                "-i",
                str(list_file),
                "-c:v",
                "libx264",
                "-pix_fmt",
                "yuv420p",
                "-preset",
                "medium",
                "-crf",
                "20",
                "-c:a",
                "aac",
                "-b:a",
                "192k",
                str(out_mp4),
            ],
            capture=True,
        )
        list_file.unlink(missing_ok=True)
        return

    # xfade chain
    durations = [audio_duration(c) for c in clips]
    inputs: list[str] = []
    for c in clips:
        inputs += ["-i", str(c)]

    # Build the filter_complex graph. xfade `offset` is the cumulative
    # duration up to the start of the transition (cumulative_so_far - crossfade).
    v_parts: list[str] = []
    a_parts: list[str] = []
    prev_v = "[0:v]"
    prev_a = "[0:a]"
    cum = durations[0]
    for i in range(1, len(clips)):
        v_out = f"[v{i}]"
        a_out = f"[a{i}]"
        offset = max(0.0, cum - crossfade)
        v_parts.append(
            f"{prev_v}[{i}:v]xfade=transition=fade:duration={crossfade}:offset={offset:.3f}{v_out}"
        )
        a_parts.append(f"{prev_a}[{i}:a]acrossfade=d={crossfade}{a_out}")
        prev_v = v_out
        prev_a = a_out
        cum = cum + durations[i] - crossfade

    v_label = prev_v
    a_label = prev_a
    filter_complex = ";".join(v_parts + a_parts)

    cmd = [
        ffmpeg,
        "-y",
        "-loglevel",
        "error",
        *inputs,
        "-filter_complex",
        filter_complex,
        "-map",
        v_label,
        "-map",
        a_label,
        "-c:v",
        "libx264",
        "-pix_fmt",
        "yuv420p",
        "-preset",
        "medium",
        "-crf",
        "20",
        "-c:a",
        "aac",
        "-b:a",
        "192k",
        str(out_mp4),
    ]
    _run(cmd, capture=True)


# ---------------------------------------------------------------------------
# Cache layout
# ---------------------------------------------------------------------------

def _cache_dir_for(pptx: Path, override: Path | None) -> Path:
    if override is not None:
        override.mkdir(parents=True, exist_ok=True)
        return override
    h = hashlib.sha1(str(pptx.resolve()).encode("utf-8")).hexdigest()[:10]
    d = Path(tempfile.gettempdir()) / f"make_video-{pptx.stem}-{h}"
    d.mkdir(parents=True, exist_ok=True)
    return d


# ---------------------------------------------------------------------------
# Driver
# ---------------------------------------------------------------------------

def build(
    pptx: Path,
    out_mp4: Path,
    *,
    tts_choice: str = "auto",
    dpi: int = 150,
    crossfade: float = 0.3,
    pre_pad: float = 1.0,
    post_pad: float = 1.5,
    cache_dir: Path | None = None,
    keep_temp: bool = False,
    lang: str = "en",
) -> dict:
    t0 = time.time()
    if not pptx.exists():
        _die(f"input .pptx not found: {pptx}")
    if pptx.suffix.lower() != ".pptx":
        _die(f"expected a .pptx file, got: {pptx.suffix}")

    cache = _cache_dir_for(pptx, cache_dir)
    _log(f"cache dir: {cache}")
    pdf_dir = cache / "pdf"
    png_dir = cache / "png"
    audio_dir = cache / "audio"
    clip_dir = cache / "clips"
    for d in (pdf_dir, png_dir, audio_dir, clip_dir):
        d.mkdir(parents=True, exist_ok=True)

    # 1. PPTX -> PDF -> PNGs
    pdf = pptx_to_pdf(pptx, pdf_dir)
    pngs = pdf_to_pngs(pdf, png_dir, dpi=dpi)

    # 2. Speaker notes
    notes = extract_notes(pptx)
    if len(notes) != len(pngs):
        _log(f"WARNING: {len(pngs)} slides vs {len(notes)} note blocks; padding the shorter one")
        while len(notes) < len(pngs):
            notes.append("(no narration provided for this slide.)")
        notes = notes[: len(pngs)]

    # 3. TTS
    tts = synth_tts(notes, audio_dir, choice=tts_choice, lang=lang)
    _log(f"TTS backend chosen: {tts.backend}")

    # 4. Per-slide subclips
    durations: list[float] = []
    clip_paths: list[Path] = []
    n = len(pngs)
    for i, (png, aud) in enumerate(zip(pngs, tts.files), start=1):
        dur = audio_duration(aud)
        # Apply pre/post pad only to first/last slide.
        pre = pre_pad if i == 1 else 0.0
        post = post_pad if i == n else 0.0
        clip = clip_dir / f"clip-{i:03d}.mp4"
        _log(f"render slide {i}/{n}: audio={dur:.2f}s pre={pre:.2f}s post={post:.2f}s")
        make_slide_clip(png, aud, clip, dur, pre_pad=pre, post_pad=post)
        durations.append(dur + pre + post)
        clip_paths.append(clip)

    # 5. Concatenate
    total_pre = sum(durations)
    _log(f"concat {n} clip(s) with crossfade={crossfade}s (sum of subclips ~{total_pre:.1f}s)")
    out_mp4.parent.mkdir(parents=True, exist_ok=True)
    concat_clips(clip_paths, out_mp4, crossfade=crossfade)

    # 6. Verify
    if not out_mp4.exists() or out_mp4.stat().st_size == 0:
        _die(f"output mp4 missing or empty: {out_mp4}")
    final_dur = audio_duration(out_mp4)
    elapsed = time.time() - t0
    _log(f"DONE: {out_mp4} ({out_mp4.stat().st_size/1e6:.1f} MB, {final_dur:.1f}s; built in {elapsed:.1f}s)")

    if not keep_temp and cache_dir is None:
        try:
            shutil.rmtree(cache)
            _log(f"cleaned cache: {cache}")
        except OSError as e:
            _log(f"could not clean cache {cache}: {e}")

    return {
        "input": str(pptx),
        "output": str(out_mp4),
        "slides": n,
        "tts_backend": tts.backend,
        "duration_seconds": final_dur,
        "size_bytes": out_mp4.stat().st_size,
        "build_seconds": elapsed,
    }


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def _parse_args(argv: list[str] | None = None) -> argparse.Namespace:
    p = argparse.ArgumentParser(
        prog="make_video.py",
        description=textwrap.dedent(__doc__ or "").strip(),
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    p.add_argument("--input", required=True, type=Path, help="path to a .pptx file")
    p.add_argument("--output", required=True, type=Path, help="path to the .mp4 to write")
    p.add_argument("--tts", choices=["auto", "gtts", "pyttsx3"], default="auto",
                   help="TTS backend (default: auto = try gTTS, fall back to pyttsx3)")
    p.add_argument("--dpi", type=int, default=150, help="pdftoppm DPI (default: 150)")
    p.add_argument("--crossfade", type=float, default=0.3, help="xfade duration in seconds (default: 0.3)")
    p.add_argument("--no-crossfade", action="store_true", help="straight concat without xfade")
    p.add_argument("--pre-pad", type=float, default=1.0, help="silence padded onto slide 1 (default: 1.0s)")
    p.add_argument("--post-pad", type=float, default=1.5, help="silence padded onto the final slide (default: 1.5s)")
    p.add_argument("--cache-dir", type=Path, default=None, help="reuse this dir for intermediate artifacts")
    p.add_argument("--keep-temp", action="store_true", help="don't delete the cache dir on exit")
    p.add_argument("--lang", default="en", help="gTTS language code (default: en)")
    p.add_argument("--json", action="store_true", help="print a JSON summary to stdout on success")
    return p.parse_args(argv)


def main(argv: list[str] | None = None) -> int:
    args = _parse_args(argv)
    crossfade = 0.0 if args.no_crossfade else max(0.0, args.crossfade)
    summary = build(
        pptx=args.input,
        out_mp4=args.output,
        tts_choice=args.tts,
        dpi=args.dpi,
        crossfade=crossfade,
        pre_pad=args.pre_pad,
        post_pad=args.post_pad,
        cache_dir=args.cache_dir,
        keep_temp=args.keep_temp,
        lang=args.lang,
    )
    if args.json:
        print(json.dumps(summary, indent=2))
    return 0


if __name__ == "__main__":
    sys.exit(main())
