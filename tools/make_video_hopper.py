#!/usr/bin/env python3
"""
make_video_hopper.py -- Hopper / A100 variant of the teaching-video pipeline.

Same CLI shape as tools/make_video.py (the local pipeline) so a SLURM batch job
can call it identically:

    python tools/make_video_hopper.py --input file.pptx --output file.mp4 [--res 4k]
                                      [--voice-sample sample.wav]

Differences from the local pipeline:
  * Neural TTS backend selection by TTS_BACKEND env var:
        xtts   -> Coqui XTTS-v2 (best quality, GPU; ~realtime on an A100)
        bark   -> suno/bark   (very expressive, slower, GPU)
        piper  -> rhasspy/piper-tts (excellent quality, CPU-only)
        gtts   -> Google Text-to-Speech (network, robotic-ish)
    Default order on Hopper:  xtts -> piper -> gtts -> pyttsx3
  * Optional voice cloning with --voice-sample sample.wav (XTTS-v2 native).
  * Higher-resolution slide rendering: --res 1080 (default) or --res 4k.
  * Optional Manim B-roll: if a slide note contains
        <!-- manim:expression -->
    a Manim animation is rendered for that slide in place of the static PNG.
    Skipped gracefully if manim is not installed.
  * ffmpeg compositing with crossfades + slight Ken-Burns zoom (matches local).

No emojis. No hard-coded secrets. GMU-cluster-specific strings live in the
SLURM scripts, not in this file.
"""
from __future__ import annotations

import argparse
import os
import re
import shlex
import shutil
import subprocess
import sys
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Optional


# ---------------------------------------------------------------------------
# Resolution presets
# ---------------------------------------------------------------------------

RES_PRESETS = {
    "720":  (1280, 720),
    "1080": (1920, 1080),
    "1440": (2560, 1440),
    "2k":   (2560, 1440),
    "4k":   (3840, 2160),
    "2160": (3840, 2160),
}


# ---------------------------------------------------------------------------
# Small helpers
# ---------------------------------------------------------------------------

def log(msg: str) -> None:
    print(f"[make_video_hopper] {msg}", flush=True)


def run(cmd: list[str], check: bool = True) -> subprocess.CompletedProcess:
    log("$ " + " ".join(shlex.quote(c) for c in cmd))
    return subprocess.run(cmd, check=check)


def have_cmd(name: str) -> bool:
    return shutil.which(name) is not None


def have_gpu() -> bool:
    """Best-effort GPU probe (Hopper compute nodes set CUDA_VISIBLE_DEVICES)."""
    if os.environ.get("CUDA_VISIBLE_DEVICES", "") not in ("", "-1"):
        return True
    if have_cmd("nvidia-smi"):
        try:
            r = subprocess.run(
                ["nvidia-smi", "-L"], capture_output=True, text=True, timeout=10
            )
            return r.returncode == 0 and "GPU" in r.stdout
        except Exception:
            return False
    return False


# ---------------------------------------------------------------------------
# Slide extraction (PPTX -> PNG + speaker notes)
# ---------------------------------------------------------------------------

@dataclass
class Slide:
    index: int
    image_path: Path
    notes: str          # speaker notes (the narration)
    manim_expr: Optional[str] = None   # populated if <!-- manim:... --> is found


def pptx_to_pngs(pptx: Path, out_dir: Path, width: int, height: int) -> list[Path]:
    """Render the pptx to one PNG per slide using libreoffice + pdftoppm.

    libreoffice writes a PDF; pdftoppm rasterizes each page to PNG at the
    requested resolution. This matches the local pipeline so a viewer cannot
    tell the rendering path apart -- only the narration changes.
    """
    out_dir.mkdir(parents=True, exist_ok=True)
    pdf_dir = out_dir / "_pdf"
    pdf_dir.mkdir(exist_ok=True)

    if not have_cmd("libreoffice") and not have_cmd("soffice"):
        raise RuntimeError(
            "libreoffice/soffice not found; cannot rasterize pptx. "
            "Load the libreoffice module or apt-install it before running."
        )
    soffice = "libreoffice" if have_cmd("libreoffice") else "soffice"

    run([soffice, "--headless", "--convert-to", "pdf",
         "--outdir", str(pdf_dir), str(pptx)])

    pdfs = sorted(pdf_dir.glob("*.pdf"))
    if not pdfs:
        raise RuntimeError(f"libreoffice produced no PDF for {pptx}")
    pdf = pdfs[0]

    # pdftoppm -r DPI gives us a width-controlled rasterization. Compute the
    # DPI needed to hit the requested pixel width on a standard 13.333"-wide
    # 16:9 slide.
    # Read the deck's ACTUAL slide width so DPI hits the target pixel width
    # exactly (pandoc's 16:9 reference is 10in wide, not 13.333in).
    try:
        from pptx import Presentation as _Prs
        slide_width_in = _Prs(str(pptx)).slide_width / 914400.0
    except Exception:
        slide_width_in = 13.333
    dpi = max(96, int(round(width / slide_width_in)))
    if not have_cmd("pdftoppm"):
        raise RuntimeError("pdftoppm not found (apt: poppler-utils).")
    run(["pdftoppm", "-png", "-r", str(dpi), str(pdf), str(out_dir / "slide")])

    pngs = sorted(out_dir.glob("slide-*.png"))
    if not pngs:
        raise RuntimeError("pdftoppm produced no PNGs")
    return pngs


def extract_notes(pptx: Path) -> list[str]:
    """Pull speaker notes per slide using python-pptx. Empty string if absent."""
    try:
        from pptx import Presentation  # python-pptx
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx not installed. pip install python-pptx"
        ) from exc

    prs = Presentation(str(pptx))
    notes: list[str] = []
    for slide in prs.slides:
        text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            text = slide.notes_slide.notes_text_frame.text or ""
        notes.append(text.strip())
    return notes


def extract_notes_from_qmd(qmd: Path) -> list[str]:
    """Return one narration string PER SLIDE (header), aligned to beamer slides.

    Walks the deck by `#`/`##` headers and grabs the `::: notes` block inside
    each header's region (empty string if a slide has none). Aligning per-header
    -- rather than concatenating all `::: notes` blocks -- keeps the narration in
    lockstep with the rendered slides even when the title slide carries no note.
    """
    text = qmd.read_text(encoding="utf-8")
    if text.startswith("---"):                  # drop YAML front matter
        text = text.split("---", 2)[-1]
    lines = text.split("\n")
    # Slide headers = `#`/`##` at column 0 that are NOT inside a ::: fenced div
    # (callout/notes blocks carry their own `##`) and NOT inside a code fence
    # (python `# comment`). Counting those as slides desyncs narration.
    depth = 0
    in_code = False
    hdr: list[int] = []
    for i, ln in enumerate(lines):
        st = ln.strip()
        if st.startswith("```"):
            in_code = not in_code
            continue
        if in_code:
            continue
        if re.match(r"^:::+", st):
            depth = max(0, depth - 1) if re.match(r"^:::+\s*$", st) else depth + 1
            continue
        if depth == 0 and re.match(r"^#{1,2}\s", ln):
            hdr.append(i)
    notes: list[str] = []
    for k, h in enumerate(hdr):
        end = hdr[k + 1] if k + 1 < len(hdr) else len(lines)
        region = "\n".join(lines[h:end])
        m = re.search(r":::+\s*notes\s*\n(.*?)\n:::", region, re.DOTALL)
        notes.append(m.group(1).strip() if m else "")
    return notes


def pdf_to_pngs(pdf: Path, out_dir: Path, width: int, height: int) -> list[Path]:
    """Rasterize a (beamer) PDF to one PNG per page at the exact target size."""
    out_dir.mkdir(parents=True, exist_ok=True)
    if not have_cmd("pdftoppm"):
        raise RuntimeError("pdftoppm not found (poppler).")
    run(["pdftoppm", "-png", "-scale-to-x", str(width), "-scale-to-y", str(height),
         str(pdf), str(out_dir / "slide")])
    pngs = sorted(out_dir.glob("slide-*.png"))
    if not pngs:
        raise RuntimeError("pdftoppm produced no PNGs from the PDF")
    return pngs


def _sentences(text: str) -> list[str]:
    return [x.strip() for x in re.split(r"(?<=[.!?])\s+", text.strip()) if x.strip()]


def _srt_ts(sec: float) -> str:
    h = int(sec // 3600); m = int((sec % 3600) // 60); s = sec % 60
    return f"{h:02d}:{m:02d}:{s:06.3f}".replace(".", ",")


def write_srt(notes: list[str], durations: list[float], path: Path) -> None:
    """Burn-in transcript: sentence segments timed proportionally per slide.

    The `::: notes` text IS the spoken transcript (symbols already spelled out),
    so on-screen captions exactly match the narration.
    """
    segs: list[tuple[float, float, str]] = []
    t = 0.0
    for note, dur in zip(notes, durations):
        sents = _sentences(note) or [note.strip() or " "]
        weights = [max(1, len(s.split())) for s in sents]
        total = sum(weights)
        st = t
        for s, w in zip(sents, weights):
            d = dur * (w / total)
            segs.append((st, st + d, s))
            st += d
        t += dur
    lines = [f"{i}\n{_srt_ts(a)} --> {_srt_ts(b)}\n{txt}\n"
             for i, (a, b, txt) in enumerate(segs, 1)]
    path.write_text("\n".join(lines), encoding="utf-8")


MANIM_MARKER = re.compile(r"<!--\s*manim:(.+?)-->", re.IGNORECASE | re.DOTALL)


def split_manim_marker(notes: str) -> tuple[str, Optional[str]]:
    """Return (clean_notes_for_TTS, manim_expression_or_None)."""
    m = MANIM_MARKER.search(notes)
    if not m:
        return notes, None
    expr = m.group(1).strip()
    clean = MANIM_MARKER.sub("", notes).strip()
    return clean, expr


# ---------------------------------------------------------------------------
# TTS backends
# ---------------------------------------------------------------------------

def _tts_xtts(text: str, out_wav: Path, voice_sample: Optional[Path]) -> bool:
    """Coqui XTTS-v2. Requires `pip install TTS` and the model cached locally."""
    try:
        from TTS.api import TTS  # type: ignore
    except ImportError:
        log("XTTS: TTS package not installed; skipping.")
        return False
    try:
        import torch  # noqa: F401
    except ImportError:
        log("XTTS: torch not installed; skipping.")
        return False
    try:
        device = "cuda" if have_gpu() else "cpu"
        log(f"XTTS: loading model on {device}")
        tts = TTS("tts_models/multilingual/multi-dataset/xtts_v2").to(device)
        kwargs = {"text": text, "file_path": str(out_wav), "language": "en"}
        if voice_sample is not None and voice_sample.exists():
            kwargs["speaker_wav"] = str(voice_sample)
        else:
            # Pick a reasonable default speaker baked into XTTS-v2.
            kwargs["speaker"] = os.environ.get("XTTS_SPEAKER", "Damien Black")
        tts.tts_to_file(**kwargs)
        return out_wav.exists() and out_wav.stat().st_size > 0
    except Exception as exc:
        log(f"XTTS failed: {exc!r}; will fall back.")
        return False


def _tts_bark(text: str, out_wav: Path) -> bool:
    try:
        from bark import SAMPLE_RATE, generate_audio, preload_models  # type: ignore
        from scipy.io.wavfile import write as wav_write  # type: ignore
    except ImportError:
        log("Bark: not installed; skipping.")
        return False
    try:
        preload_models()
        audio = generate_audio(text)
        wav_write(str(out_wav), SAMPLE_RATE, audio)
        return out_wav.exists()
    except Exception as exc:
        log(f"Bark failed: {exc!r}; will fall back.")
        return False


# Module-level cache so the Chatterbox model loads ONCE per process, not per
# slide (loading is the expensive part; synthesis is faster-than-realtime on A100).
_CHATTERBOX_MODEL = None


def _chatterbox_chunks(text: str, max_chars: int = 300) -> list[str]:
    """Split narration into sentence groups Chatterbox can synthesize in one pass."""
    sentences = re.split(r"(?<=[.!?])\s+", text.strip())
    chunks, cur = [], ""
    for s in sentences:
        if not s:
            continue
        if len(cur) + len(s) + 1 <= max_chars:
            cur = f"{cur} {s}".strip()
        else:
            if cur:
                chunks.append(cur)
            cur = s
    if cur:
        chunks.append(cur)
    return chunks or [text.strip()]


def _tts_chatterbox(text: str, out_wav: Path, voice_sample: Optional[Path]) -> bool:
    """Chatterbox (Resemble AI) neural TTS -- MIT (code+weights), commercial-safe.

    Zero-shot voice cloning from a short reference WAV via `audio_prompt_path`;
    falls back to the model's built-in voice if no sample is given. Long notes are
    synthesized in sentence-grouped chunks and concatenated. Every clip carries
    Chatterbox's inaudible PerTh watermark (harmless for teaching).
    """
    global _CHATTERBOX_MODEL
    try:
        import torch  # noqa: F401
        import torchaudio  # type: ignore
        from chatterbox.tts import ChatterboxTTS  # type: ignore
    except ImportError:
        log("Chatterbox: package/torch/torchaudio not installed; skipping.")
        return False
    try:
        device = "cuda" if have_gpu() else "cpu"
        if _CHATTERBOX_MODEL is None:
            log(f"Chatterbox: loading model on {device} (one-time)")
            _CHATTERBOX_MODEL = ChatterboxTTS.from_pretrained(device=device)
        model = _CHATTERBOX_MODEL

        prompt = None
        if voice_sample is not None and voice_sample.exists():
            prompt = str(voice_sample)
        gen_kwargs = {}
        if prompt:
            gen_kwargs["audio_prompt_path"] = prompt
        # Calmer, lecture-appropriate delivery; override via env if desired.
        gen_kwargs["exaggeration"] = float(os.environ.get("CHATTERBOX_EXAGGERATION", "0.4"))
        gen_kwargs["cfg_weight"] = float(os.environ.get("CHATTERBOX_CFG", "0.5"))

        waves = []
        for chunk in _chatterbox_chunks(text):
            wav = model.generate(chunk, **gen_kwargs)
            waves.append(wav)
        audio = waves[0] if len(waves) == 1 else torch.cat(waves, dim=-1)
        torchaudio.save(str(out_wav), audio.cpu(), model.sr)
        return out_wav.exists() and out_wav.stat().st_size > 0
    except Exception as exc:
        log(f"Chatterbox failed: {exc!r}; will fall back.")
        return False


def _tts_piper(text: str, out_wav: Path) -> bool:
    """Piper CLI TTS -- CPU-only, very good quality."""
    if not have_cmd("piper"):
        log("Piper: binary not found; skipping.")
        return False
    voice = os.environ.get("PIPER_VOICE", "en_US-amy-medium.onnx")
    if not Path(voice).exists():
        log(f"Piper: voice file {voice} not found; skipping.")
        return False
    try:
        proc = subprocess.run(
            ["piper", "--model", voice, "--output_file", str(out_wav)],
            input=text, text=True, check=True,
        )
        return proc.returncode == 0 and out_wav.exists()
    except Exception as exc:
        log(f"Piper failed: {exc!r}; will fall back.")
        return False


def _tts_gtts(text: str, out_wav: Path) -> bool:
    try:
        from gtts import gTTS  # type: ignore
    except ImportError:
        log("gTTS: not installed; skipping.")
        return False
    try:
        mp3 = out_wav.with_suffix(".mp3")
        gTTS(text=text, lang="en").save(str(mp3))
        # transcode mp3 -> wav so the rest of the pipeline is format-agnostic
        run(["ffmpeg", "-y", "-loglevel", "error",
             "-i", str(mp3), str(out_wav)])
        mp3.unlink(missing_ok=True)
        return out_wav.exists()
    except Exception as exc:
        log(f"gTTS failed: {exc!r}; will fall back.")
        return False


def _tts_pyttsx3(text: str, out_wav: Path) -> bool:
    try:
        import pyttsx3  # type: ignore
    except ImportError:
        log("pyttsx3: not installed; nothing left to fall back to.")
        return False
    try:
        engine = pyttsx3.init()
        engine.save_to_file(text, str(out_wav))
        engine.runAndWait()
        return out_wav.exists()
    except Exception as exc:
        log(f"pyttsx3 failed: {exc!r}.")
        return False


def synthesize_narration(text: str, out_wav: Path,
                         voice_sample: Optional[Path]) -> Path:
    """Synthesize `text` to `out_wav` using the requested/auto backend.

    Order: env TTS_BACKEND if set, then xtts (if GPU) -> piper -> gtts -> pyttsx3.
    """
    if not text.strip():
        # Generate ~0.7s of silence so the slide still has a duration.
        run(["ffmpeg", "-y", "-loglevel", "error",
             "-f", "lavfi", "-i", "anullsrc=r=24000:cl=mono",
             "-t", "0.7", str(out_wav)])
        return out_wav

    requested = os.environ.get("TTS_BACKEND", "").lower().strip()
    # Chatterbox first: MIT-licensed (code+weights), commercial-safe, voice cloning.
    auto_order = ["chatterbox", "xtts", "piper", "gtts", "pyttsx3"]
    if not have_gpu():
        # GPU-only neural backends are slow on CPU; demote them below piper.
        for b in ("chatterbox", "xtts"):
            if b in auto_order:
                auto_order.remove(b)
                auto_order.insert(2, b)
    order = [requested] + [b for b in auto_order if b != requested] if requested else auto_order

    for backend in order:
        if not backend:
            continue
        log(f"TTS: trying backend={backend}")
        ok = False
        if backend == "chatterbox":
            ok = _tts_chatterbox(text, out_wav, voice_sample)
        elif backend == "xtts":
            ok = _tts_xtts(text, out_wav, voice_sample)
        elif backend == "bark":
            ok = _tts_bark(text, out_wav)
        elif backend == "piper":
            ok = _tts_piper(text, out_wav)
        elif backend == "gtts":
            ok = _tts_gtts(text, out_wav)
        elif backend == "pyttsx3":
            ok = _tts_pyttsx3(text, out_wav)
        else:
            log(f"Unknown backend {backend!r}, skipping.")
            continue
        if ok:
            return out_wav

    raise RuntimeError(
        "All TTS backends failed. Install at least one of: TTS, piper, gtts, pyttsx3."
    )


# ---------------------------------------------------------------------------
# Optional Manim B-roll
# ---------------------------------------------------------------------------

def render_manim(expr: str, out_png: Path, width: int, height: int) -> bool:
    """Render a Manim animation for this slide. Returns True on success.

    `expr` should be a self-contained Manim Scene name or a tiny snippet
    (e.g. "PlotCAPM"). To keep this safe, we only accept identifier-like
    names and look them up in tools/manim_scenes.py. Free-form expressions
    are NOT exec'd here.
    """
    if not re.fullmatch(r"[A-Za-z_][A-Za-z0-9_]*", expr or ""):
        log(f"Manim expr {expr!r} is not a bare Scene name; skipping for safety.")
        return False
    if not have_cmd("manim"):
        log("Manim: not installed; skipping B-roll.")
        return False
    scenes_py = Path(__file__).parent / "manim_scenes.py"
    if not scenes_py.exists():
        log(f"Manim: {scenes_py} not found; skipping.")
        return False
    with tempfile.TemporaryDirectory() as td:
        td_p = Path(td)
        try:
            run(["manim", "-ql", "-o", "broll", "--media_dir", str(td_p),
                 str(scenes_py), expr])
        except subprocess.CalledProcessError as exc:
            log(f"Manim render failed: {exc}; skipping.")
            return False
        mp4s = list(td_p.rglob("*.mp4"))
        if not mp4s:
            log("Manim: no mp4 produced; skipping.")
            return False
        # Grab first/last frame as a PNG fallback so the rest of the pipeline
        # (which composites stills + audio) still works. For full B-roll
        # treatment we would splice the mp4 into the timeline; that is a
        # future extension.
        run(["ffmpeg", "-y", "-loglevel", "error",
             "-sseof", "-0.1", "-i", str(mp4s[0]),
             "-vframes", "1", "-vf", f"scale={width}:{height}",
             str(out_png)])
        return out_png.exists()


# ---------------------------------------------------------------------------
# ffmpeg compositing (matches local pipeline: per-slide MP4 with Ken-Burns,
# then concat with crossfades)
# ---------------------------------------------------------------------------

def audio_duration(path: Path) -> float:
    r = subprocess.run(
        ["ffprobe", "-v", "error", "-show_entries", "format=duration",
         "-of", "default=noprint_wrappers=1:nokey=1", str(path)],
        capture_output=True, text=True, check=True,
    )
    return float(r.stdout.strip() or "1.0")


def build_clip(png: Path, wav: Path, out_mp4: Path,
               width: int, height: int) -> None:
    """One slide -> one MP4 with a slow Ken-Burns zoom matched to audio length."""
    dur = max(1.5, audio_duration(wav) + 0.4)
    # zoompan zooms slowly from 1.0 to ~1.06 over the clip
    fps = 30
    frames = int(dur * fps)
    zoom_expr = f"zoom='min(zoom+0.0005,1.06)':d={frames}:s={width}x{height}:fps={fps}"
    vf = (
        f"scale={width*2}:{height*2},"
        f"zoompan={zoom_expr},"
        f"format=yuv420p"
    )
    run([
        "ffmpeg", "-y", "-loglevel", "error",
        "-loop", "1", "-i", str(png),
        "-i", str(wav),
        "-c:v", "libx264", "-preset", "medium", "-crf", "20",
        "-c:a", "aac", "-b:a", "192k",
        "-vf", vf,
        "-t", f"{dur:.3f}",
        "-pix_fmt", "yuv420p", "-shortest",
        str(out_mp4),
    ])


def concat_with_crossfades(clips: list[Path], out_mp4: Path,
                           width: int, height: int) -> None:
    """Concatenate per-slide MP4s with a 0.4s crossfade between each pair.

    For N clips we emit an ffmpeg filter_complex with N-1 xfade/acrossfade
    pairs. For small N this is fine; for very large N a simple concat demuxer
    (no crossfade) would be cheaper -- but our weekly decks are ~20-40 slides.
    """
    if not clips:
        raise RuntimeError("No clips to concatenate")
    if len(clips) == 1:
        shutil.copy2(clips[0], out_mp4)
        return

    # Concat demuxer with re-encode. Chained xfade offsets are fragile for N>2
    # (they silently truncate the output); a plain concat gives exact, correct
    # total duration with clean hard cuts between slides -- right for lectures.
    list_file = out_mp4.parent / "_concat_list.txt"
    list_file.write_text("".join(f"file '{Path(c).resolve()}'\n" for c in clips))
    run([
        "ffmpeg", "-y", "-loglevel", "error",
        "-f", "concat", "-safe", "0", "-i", str(list_file),
        "-c:v", "libx264", "-preset", "medium", "-crf", "20",
        "-c:a", "aac", "-b:a", "192k",
        "-pix_fmt", "yuv420p",
        str(out_mp4),
    ])


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> int:
    ap = argparse.ArgumentParser(
        description="Hopper-side teaching-video render: PPTX -> narrated MP4."
    )
    ap.add_argument("--input", required=True, type=Path,
                    help="Slides: a beamer .pdf (preferred) or legacy .pptx.")
    ap.add_argument("--notes-from", type=Path, default=None,
                    help="Deck .qmd to read narration from (required for .pdf input).")
    ap.add_argument("--output", required=True, type=Path,
                    help="Destination .mp4.")
    ap.add_argument("--res", default="1080",
                    choices=sorted(RES_PRESETS.keys()),
                    help="Output resolution preset (default 1080).")
    ap.add_argument("--voice-sample", type=Path, default=None,
                    help="Optional 20-30s WAV for XTTS voice cloning.")
    ap.add_argument("--keep-work", action="store_true",
                    help="Keep the intermediate working directory.")
    args = ap.parse_args()

    pptx: Path = args.input
    out_mp4: Path = args.output
    if not pptx.exists():
        log(f"ERROR: input not found: {pptx}")
        return 2
    out_mp4.parent.mkdir(parents=True, exist_ok=True)

    width, height = RES_PRESETS[args.res]
    log(f"Resolution: {width}x{height} (--res {args.res})")
    log(f"GPU available: {have_gpu()}")

    work = Path(tempfile.mkdtemp(prefix="mvh_"))
    log(f"Working dir: {work}")
    try:
        # 1) Slides -> per-slide PNGs; narration list aligned 1:1 to slides.
        if pptx.suffix.lower() == ".pdf":
            pngs = pdf_to_pngs(pptx, work / "slides", width, height)
            notes_src = args.notes_from or pptx.with_suffix(".qmd")
            notes_list = extract_notes_from_qmd(Path(notes_src))
            log(f"beamer path: {len(pngs)} slides, {len(notes_list)} note blocks "
                f"(notes from {notes_src})")
        else:
            pngs = pptx_to_pngs(pptx, work / "slides", width, height)
            notes_list = extract_notes(pptx)
        if len(notes_list) < len(pngs):
            notes_list += [""] * (len(pngs) - len(notes_list))
        elif len(notes_list) > len(pngs):
            log(f"WARN: {len(notes_list)} notes > {len(pngs)} slides; truncating")
            notes_list = notes_list[:len(pngs)]

        # 2) Per slide: TTS narration (+ optional manim B-roll image)
        clips: list[Path] = []
        cap_notes: list[str] = []
        cap_durs: list[float] = []
        for i, png in enumerate(pngs):
            raw_notes = notes_list[i] if i < len(notes_list) else ""
            clean_notes, manim_expr = split_manim_marker(raw_notes)
            if manim_expr:
                broll_png = work / f"manim-{i+1:03d}.png"
                if render_manim(manim_expr, broll_png, width, height):
                    png = broll_png  # replace the slide still with the manim frame
            wav = work / f"slide-{i+1:03d}.wav"
            synthesize_narration(clean_notes or " ", wav, args.voice_sample)
            clip = work / f"clip-{i+1:03d}.mp4"
            # ensure the wav sits next to the clip so concat_with_crossfades
            # can locate it for duration probing
            shutil.copy2(wav, clip.with_suffix(".wav"))
            build_clip(png, wav, clip, width, height)
            clips.append(clip)
            cap_notes.append(clean_notes)
            cap_durs.append(audio_duration(wav))

        # 3) Concatenate, and emit a transcript .srt next to the output.
        concat_with_crossfades(clips, out_mp4, width, height)
        try:
            write_srt(cap_notes, cap_durs, out_mp4.with_suffix(".srt"))
            log(f"Wrote transcript {out_mp4.with_suffix('.srt')}")
        except Exception as exc:
            log(f"SRT generation failed (non-fatal): {exc!r}")
        log(f"Wrote {out_mp4} ({out_mp4.stat().st_size/1e6:.1f} MB)")
        return 0
    finally:
        if not args.keep_work:
            shutil.rmtree(work, ignore_errors=True)


if __name__ == "__main__":
    sys.exit(main())
