#!/usr/bin/env python3
"""
build_reference_pptx.py — generate the GMU-branded Quarto/Pandoc reference deck.

Output: book/decks/templates/assip-reference.pptx

Quarto/Pandoc cannot theme PPTX with SCSS/CSS; the ONLY supported mechanism is a
reference `.pptx` whose Slide Master + theme Pandoc copies into every rendered
deck (see _template.qmd `reference-doc:`). This script starts from python-pptx's
default template (whose layout NAMES already match what Pandoc expects —
"Title Slide", "Title and Content", "Section Header", ...) and rewrites:

  * the theme color scheme  -> GMU April-2024 rebrand palette
  * the theme font scheme   -> Poppins (headings) / Open Sans (body), Office-safe
  * a footer wordmark + a Mason-green accent bar on the slide master

The real interlocking-GM logo is NetID/Bynder-gated and is NOT bundled; drop it
into assets/brand/ and re-run with --logo to stamp it on the masters.

Verified GMU hex values (gmu.edu brand guidelines, April-2024 rebrand):
  Mason Green #005239  ·  Mason Gold #FFC733  ·  Logo Black #333333
  accents (charts): Red #CC4824 · Teal #008285 · Navy #004F71 · Gray #727579
Do NOT use the legacy #006633 / #FFCC33.
"""
from __future__ import annotations
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree
from PIL import Image

R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# ---- GMU brand (verified) ---------------------------------------------------
GMU = {
    "green":  "005239",   # Mason Green  (primary)
    "gold":   "FFC733",   # Mason Gold   (primary)
    "black":  "333333",   # Logo Black   (text)
    "white":  "FFFFFF",
    "paper":  "F4F4F2",   # subtle off-white background
    "red":    "CC4824",   # accent (charts)
    "teal":   "008285",   # accent
    "navy":   "004F71",   # accent
    "gray":   "727579",   # accent
}

# Map the Office theme slots (dk1/lt1/dk2/lt2/accent1..6/hlink/folHlink) -> GMU.
# accent1 = Green and accent2 = Gold so Pandoc's default styling reads on-brand.
THEME_COLORS = {
    "dk1": GMU["black"], "lt1": GMU["white"],
    "dk2": GMU["green"], "lt2": GMU["paper"],
    "accent1": GMU["green"], "accent2": GMU["gold"], "accent3": GMU["red"],
    "accent4": GMU["teal"],  "accent5": GMU["navy"], "accent6": GMU["gray"],
    "hlink": GMU["navy"], "folHlink": GMU["teal"],
}
HEAD_FONT = "Poppins"     # display/headings (Office fallback: Arial/Calibri)
BODY_FONT = "Open Sans"   # body            (Office fallback: Calibri/Arial)

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS = {"a": A}


def _qn(tag: str) -> str:
    return f"{{{A}}}{tag}"


def _set_srgb(parent):
    """Replace whatever color child a theme slot has with a plain <a:srgbClr>."""
    for child in list(parent):
        parent.remove(child)
    srgb = etree.SubElement(parent, _qn("srgbClr"))
    return srgb


def recolor_and_refont_theme(theme_xml: bytes) -> bytes:
    root = etree.fromstring(theme_xml)
    clr = root.find(f".//a:themeElements/a:clrScheme", NS)
    for slot, hexval in THEME_COLORS.items():
        node = clr.find(f"a:{slot}", NS)
        if node is None:
            continue
        srgb = _set_srgb(node)
        srgb.set("val", hexval.upper())
    font = root.find(".//a:themeElements/a:fontScheme", NS)
    for role, typeface in (("majorFont", HEAD_FONT), ("minorFont", BODY_FONT)):
        latin = font.find(f"a:{role}/a:latin", NS)
        if latin is not None:
            latin.set("typeface", typeface)
    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def stamp_master(prs, logo: Path | None):
    """Add a Mason-green accent bar + wordmark footer to every master via raw XML.

    python-pptx's MasterShapes does not expose add_shape/add_textbox, so we
    append the shapes' DrawingML directly to the master spTree.
    """
    EMU = 914400
    W, H = prs.slide_width, prs.slide_height
    bar_h = int(0.22 * EMU)
    bar = (
        f'<p:sp xmlns:p="{P_NS}" xmlns:a="{A}"><p:nvSpPr>'
        f'<p:cNvPr id="900" name="GMU Accent Bar"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="0" y="{H - bar_h}"/><a:ext cx="{W}" cy="{bar_h}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="{GMU["green"]}"/></a:solidFill>'
        f'<a:ln><a:noFill/></a:ln></p:spPr>'
        f'<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
    )
    wm_x, wm_cx, wm_cy = int(0.4 * EMU), int(7 * EMU), int(0.3 * EMU)
    wm = (
        f'<p:sp xmlns:p="{P_NS}" xmlns:a="{A}"><p:nvSpPr>'
        f'<p:cNvPr id="901" name="GMU Wordmark"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{wm_x}" y="{H - int(0.55 * EMU)}"/>'
        f'<a:ext cx="{wm_cx}" cy="{wm_cy}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/></p:spPr>'
        f'<p:txBody><a:bodyPr wrap="none"/><a:lstStyle/><a:p><a:r>'
        f'<a:rPr lang="en-US" sz="900" b="1"><a:solidFill><a:srgbClr val="{GMU["green"]}"/></a:solidFill>'
        f'<a:latin typeface="{BODY_FONT}"/></a:rPr>'
        f'<a:t>NextGen Empirical Finance Research Camp  ·  Prof. Lei Gao</a:t>'
        f'</a:r></a:p></p:txBody></p:sp>'
    )
    for master in prs.slide_masters:
        spTree = master.shapes._spTree
        spTree.append(etree.fromstring(bar))
        spTree.append(etree.fromstring(wm))
        if logo and logo.exists():
            # MasterShapes has no add_picture; embed via the image-part API + raw <p:pic>.
            try:
                with Image.open(str(logo)) as im:
                    aspect = im.width / im.height
                cy = int(0.55 * EMU); cx = int(cy * aspect)
                left = int(W - cx - 0.3 * EMU); top = int(0.28 * EMU)
                _, rId = master.part.get_or_add_image_part(str(logo))
                pic = (
                    f'<p:pic xmlns:p="{P_NS}" xmlns:a="{A}" xmlns:r="{R_NS}"><p:nvPicPr>'
                    f'<p:cNvPr id="902" name="GMU Logo"/>'
                    f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr><p:nvPr/></p:nvPicPr>'
                    f'<p:blipFill><a:blip r:embed="{rId}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
                    f'<p:spPr><a:xfrm><a:off x="{left}" y="{top}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>'
                    f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>'
                )
                spTree.append(etree.fromstring(pic))
            except Exception as exc:
                print(f"WARN: logo stamp failed ({exc!r}); template built without logo.")


def build(out: Path, logo: Path | None):
    prs = Presentation()  # default template: Pandoc-compatible layout names
    # 16:9 widescreen (13.333 x 7.5 in) so rendered video is native 1080p, no pillarbox
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    # rewrite every theme part in the package
    n_themes = 0
    for part in prs.part.package.iter_parts():
        if "/theme/" in str(part.partname) and str(part.partname).endswith(".xml"):
            part._blob = recolor_and_refont_theme(part.blob)
            n_themes += 1
    stamp_master(prs, logo)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return n_themes


def verify(out: Path):
    prs = Presentation(str(out))
    for part in prs.part.package.iter_parts():
        if "/theme/" in str(part.partname) and str(part.partname).endswith(".xml"):
            root = etree.fromstring(part.blob)
            a1 = root.find(".//a:clrScheme/a:accent1/a:srgbClr", NS)
            a2 = root.find(".//a:clrScheme/a:accent2/a:srgbClr", NS)
            maj = root.find(".//a:fontScheme/a:majorFont/a:latin", NS)
            return (a1.get("val") if a1 is not None else None,
                    a2.get("val") if a2 is not None else None,
                    maj.get("typeface") if maj is not None else None)
    return None


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=str(Path(__file__).with_name("nextgen-reference.pptx")))
    ap.add_argument("--logo", default=None, help="path to GMU logo PNG (optional)")
    args = ap.parse_args()
    out = Path(args.out)
    logo = Path(args.logo) if args.logo else None
    n = build(out, logo)
    a1, a2, maj = verify(out)
    print(f"wrote {out}  ({n} theme part(s) rebranded)")
    print(f"verify -> accent1(green)={a1}  accent2(gold)={a2}  headingFont={maj}")
